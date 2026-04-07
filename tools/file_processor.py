"""
File processor — handles all supported upload types.

Processing modes:
  auto   — images → vision, documents → text extraction  (default)
  text   — all files → text extraction only (images get a basic note)
  vision — all files → Gemma4 vision (PDFs rendered to images per-page)
  both   — text extraction + vision, results combined

Image types  → Gemma4 vision on 192.168.1.220
PDF          → pymupdf4llm markdown  /  per-page vision render
DOCX/DOC     → python-docx plain text
PPTX         → python-pptx slide text
XLSX/XLS     → openpyxl CSV-style text
TXT/MD/JSON/YAML/code → direct UTF-8 decode
"""
import asyncio
import base64
import io
import json
import mimetypes
import os
import re
from pathlib import Path
from typing import Optional

import chardet
import httpx

# ── Constants ──────────────────────────────────────────────────────────────────

VISION_URL   = os.getenv("OLLAMA_REMOTE_URL", "http://192.168.1.220:11434")
VISION_MODEL = os.getenv("VISION_MODEL",      "gemma4:latest")

# Processing modes
MODE_AUTO   = "auto"    # smart: images→vision, docs→text
MODE_TEXT   = "text"    # force text extraction for everything
MODE_VISION = "vision"  # force Gemma4 vision for everything
MODE_BOTH   = "both"    # text extraction + vision, combined

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif", ".svg"}
PDF_EXTS   = {".pdf"}
WORD_EXTS  = {".docx", ".doc"}
PPT_EXTS   = {".pptx", ".ppt"}
XL_EXTS    = {".xlsx", ".xls", ".ods"}
TEXT_EXTS  = {
    ".txt", ".md", ".markdown", ".rst",
    ".json", ".jsonl", ".ndjson",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".xml", ".html", ".htm", ".svg",
    ".csv", ".tsv",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".mjs", ".cjs",
    ".rs", ".go", ".java", ".kt", ".swift", ".cpp", ".c", ".h", ".hpp",
    ".cs", ".rb", ".php", ".lua", ".r", ".jl", ".ex", ".exs", ".elm",
    ".sh", ".bash", ".zsh", ".fish", ".ps1", ".bat", ".cmd",
    ".sql", ".graphql", ".proto",
    ".css", ".scss", ".sass", ".less",
    ".log", ".env", ".gitignore", ".dockerignore", ".editorconfig",
    ".Dockerfile", ".makefile", ".mk",
    ".ipynb",
}
MAX_TEXT_BYTES   = 5 * 1024 * 1024   # 5 MB text limit per file
MAX_IMAGE_BYTES  = 20 * 1024 * 1024  # 20 MB image limit
PDF_VISION_DPI   = 150               # DPI for PDF→image render
PDF_VISION_MAX_PAGES = 10            # Max pages to send to vision (cost/time)


class ProcessingError(Exception):
    pass


# ── Dispatcher ──────────────────────────────────────────────────────────────────

async def process_file(filename: str, data: bytes, mode: str = MODE_AUTO) -> dict:
    """
    Main entry point.

    Returns:
    {
        "filename": str,
        "ext": str,
        "type": str,              # "image" | "pdf" | "word" | "pptx" | "excel" | "text"
        "mode_used": str,         # which mode was actually applied
        "text": str,              # extracted text for RAG
        "text_extracted": str,    # text-only extraction (when mode=both)
        "vision_description": str | None,
        "pages": int | None,
        "size_bytes": int,
        "truncated": bool,
    }
    """
    ext = Path(filename).suffix.lower()
    if not ext and "." not in filename:
        ext = _sniff_ext(data)

    size = len(data)
    result = {
        "filename":          filename,
        "ext":               ext,
        "size_bytes":        size,
        "truncated":         False,
        "vision_description": None,
        "text_extracted":    None,
        "pages":             None,
        "mode_used":         mode,
    }

    is_img = ext in IMAGE_EXTS or _is_image(data)

    # ── IMAGE ────────────────────────────────────────────────────────────────
    if is_img:
        result["type"] = "image"
        if mode == MODE_TEXT:
            # Text-only: no vision model — return a note instead
            result["text"] = f"[Image file: {filename}. Vision analysis disabled (text-only mode).]"
            result["mode_used"] = MODE_TEXT
        else:
            # auto / vision / both → always run vision for images
            desc = await _process_image(filename, data)
            result["text"] = desc
            result["vision_description"] = desc
            result["mode_used"] = MODE_VISION

    # ── PDF ──────────────────────────────────────────────────────────────────
    elif ext in PDF_EXTS:
        result["type"] = "pdf"
        if mode == MODE_VISION:
            vision_text, pages = await _process_pdf_vision(data)
            result["text"] = vision_text
            result["vision_description"] = vision_text
            result["pages"] = pages
            result["mode_used"] = MODE_VISION
        elif mode == MODE_BOTH:
            text, pages = _process_pdf(data)
            vision_text, _ = await _process_pdf_vision(data)
            combined = f"## Text Extraction\n\n{text}\n\n---\n\n## Vision Analysis\n\n{vision_text}"
            result["text"] = combined
            result["text_extracted"] = text
            result["vision_description"] = vision_text
            result["pages"] = pages
            result["mode_used"] = MODE_BOTH
        else:
            # auto or text
            text, pages = _process_pdf(data)
            result["text"] = text
            result["pages"] = pages
            result["mode_used"] = MODE_TEXT

    # ── WORD ─────────────────────────────────────────────────────────────────
    elif ext in WORD_EXTS:
        result["type"] = "word"
        if mode == MODE_VISION:
            vision_text = await _process_doc_vision(filename, data)
            result["text"] = vision_text
            result["vision_description"] = vision_text
            result["mode_used"] = MODE_VISION
        elif mode == MODE_BOTH:
            text = _process_docx(data, ext)
            vision_text = await _process_doc_vision(filename, data)
            combined = f"## Text Extraction\n\n{text}\n\n---\n\n## Vision Analysis\n\n{vision_text}"
            result["text"] = combined
            result["text_extracted"] = text
            result["vision_description"] = vision_text
            result["mode_used"] = MODE_BOTH
        else:
            result["text"] = _process_docx(data, ext)
            result["mode_used"] = MODE_TEXT

    # ── PPTX ─────────────────────────────────────────────────────────────────
    elif ext in PPT_EXTS:
        result["type"] = "pptx"
        if mode == MODE_VISION:
            vision_text = await _process_doc_vision(filename, data)
            result["text"] = vision_text
            result["vision_description"] = vision_text
            result["mode_used"] = MODE_VISION
        elif mode == MODE_BOTH:
            text, pages = _process_pptx(data)
            vision_text = await _process_doc_vision(filename, data)
            combined = f"## Text Extraction\n\n{text}\n\n---\n\n## Vision Analysis\n\n{vision_text}"
            result["text"] = combined
            result["text_extracted"] = text
            result["vision_description"] = vision_text
            result["pages"] = pages
            result["mode_used"] = MODE_BOTH
        else:
            text, pages = _process_pptx(data)
            result["text"] = text
            result["pages"] = pages
            result["mode_used"] = MODE_TEXT

    # ── EXCEL ─────────────────────────────────────────────────────────────────
    elif ext in XL_EXTS:
        result["type"] = "excel"
        if mode == MODE_VISION:
            vision_text = await _process_doc_vision(filename, data)
            result["text"] = vision_text
            result["vision_description"] = vision_text
            result["mode_used"] = MODE_VISION
        elif mode == MODE_BOTH:
            text, pages = _process_excel(data, ext)
            vision_text = await _process_doc_vision(filename, data)
            combined = f"## Spreadsheet Data\n\n{text}\n\n---\n\n## Vision Analysis\n\n{vision_text}"
            result["text"] = combined
            result["text_extracted"] = text
            result["vision_description"] = vision_text
            result["pages"] = pages
            result["mode_used"] = MODE_BOTH
        else:
            text, pages = _process_excel(data, ext)
            result["text"] = text
            result["pages"] = pages
            result["mode_used"] = MODE_TEXT

    # ── TEXT / CODE ──────────────────────────────────────────────────────────
    else:
        result["type"] = "text"
        if mode == MODE_VISION:
            # Wrap raw text in a message for Gemma4 — useful for code review
            raw = _decode_text(data, filename)
            vision_text = await _analyze_text_with_vision(filename, raw)
            result["text"] = vision_text
            result["vision_description"] = vision_text
            result["text_extracted"] = raw
            result["mode_used"] = MODE_VISION
        elif mode == MODE_BOTH:
            raw = _decode_text(data, filename)
            vision_text = await _analyze_text_with_vision(filename, raw)
            combined = f"## Raw Content\n\n{raw}\n\n---\n\n## AI Analysis\n\n{vision_text}"
            result["text"] = combined
            result["text_extracted"] = raw
            result["vision_description"] = vision_text
            result["mode_used"] = MODE_BOTH
        else:
            result["text"] = _decode_text(data, filename)
            result["mode_used"] = MODE_TEXT

    # Truncate oversized text
    text = result.get("text", "")
    if len(text) > MAX_TEXT_BYTES:
        result["text"] = text[:MAX_TEXT_BYTES]
        result["truncated"] = True

    return result


# ── Image → Gemma4 vision ──────────────────────────────────────────────────────

async def _process_image(filename: str, data: bytes) -> str:
    if len(data) > MAX_IMAGE_BYTES:
        raise ProcessingError(f"Image too large: {len(data) / 1e6:.1f} MB (max 20 MB)")

    b64 = base64.b64encode(data).decode()
    ext = Path(filename).suffix.lower().lstrip(".")
    mime_map = {"jpg": "jpeg", "tif": "tiff"}
    mime_ext = mime_map.get(ext, ext) or "jpeg"

    prompt = (
        "Analyze this image thoroughly. Do ALL of the following:\n"
        "1. Extract ALL visible text exactly as written (OCR).\n"
        "2. Describe the image content, layout, and context.\n"
        "3. If it contains data, tables, charts, or diagrams — describe their content and values.\n"
        "4. Note any code, URLs, emails, or structured data visible.\n"
        "5. Describe colors, style, and visual design if relevant.\n"
        "Format your response in clear sections."
    )

    try:
        async with httpx.AsyncClient(timeout=120) as client:
            r = await client.post(
                f"{VISION_URL}/api/chat",
                json={
                    "model": VISION_MODEL,
                    "messages": [{
                        "role": "user",
                        "content": prompt,
                        "images": [b64],
                    }],
                    "stream": False,
                },
            )
            data_out = r.json()
            return data_out.get("message", {}).get("content", "").strip()
    except Exception as e:
        raise ProcessingError(f"Vision analysis failed: {e}")


async def _process_pdf_vision(data: bytes) -> tuple[str, int]:
    """Render each PDF page to an image and run Gemma4 vision on each."""
    try:
        import pymupdf
        doc = pymupdf.open(stream=data, filetype="pdf")
        pages = doc.page_count
        limit = min(pages, PDF_VISION_MAX_PAGES)
        parts = []

        prompt = (
            "Analyze this document page image completely:\n"
            "1. Extract ALL text exactly as written.\n"
            "2. Describe tables, charts, diagrams, and their data.\n"
            "3. Note any images, logos, or visual elements.\n"
            "4. Preserve document structure (headings, bullets, sections).\n"
            "Be thorough and accurate."
        )

        for i in range(limit):
            page = doc[i]
            pix = page.get_pixmap(dpi=PDF_VISION_DPI)
            img_bytes = pix.tobytes("png")
            b64 = base64.b64encode(img_bytes).decode()

            try:
                async with httpx.AsyncClient(timeout=120) as client:
                    r = await client.post(
                        f"{VISION_URL}/api/chat",
                        json={
                            "model": VISION_MODEL,
                            "messages": [{"role": "user", "content": prompt, "images": [b64]}],
                            "stream": False,
                        },
                    )
                    content = r.json().get("message", {}).get("content", "").strip()
                    if content:
                        parts.append(f"## Page {i+1}\n\n{content}")
            except Exception as e:
                parts.append(f"## Page {i+1}\n\n[Vision failed: {e}]")

        doc.close()
        if pages > limit:
            parts.append(f"\n\n*Note: Only first {limit} of {pages} pages analyzed via vision.*")
        return "\n\n---\n\n".join(parts), pages

    except Exception as e:
        raise ProcessingError(f"PDF vision failed: {e}")


async def _process_doc_vision(filename: str, data: bytes) -> str:
    """Send a document file directly to Gemma4 for vision analysis (as raw bytes)."""
    # For non-image docs, we pass the raw content as a text prompt to the LLM
    # (vision on arbitrary binary files isn't supported — do a text-based analysis instead)
    try:
        # Extract raw text best-effort, then ask Gemma4 to analyze it
        ext = Path(filename).suffix.lower()
        try:
            if ext in PDF_EXTS:
                raw, _ = _process_pdf(data)
            elif ext in WORD_EXTS:
                raw = _process_docx(data, ext)
            elif ext in PPT_EXTS:
                raw, _ = _process_pptx(data)
            elif ext in XL_EXTS:
                raw, _ = _process_excel(data, ext)
            else:
                raw = _decode_text(data, filename)
        except Exception:
            raw = data.decode("utf-8", errors="replace")

        return await _analyze_text_with_vision(filename, raw)
    except Exception as e:
        raise ProcessingError(f"Document vision analysis failed: {e}")


async def _analyze_text_with_vision(filename: str, text: str) -> str:
    """Send extracted text to Gemma4 for deep analysis / enrichment."""
    ext = Path(filename).suffix.lower()
    truncated = text[:12000] if len(text) > 12000 else text  # keep within context

    type_hint = {
        ".py": "Python code", ".js": "JavaScript code", ".ts": "TypeScript code",
        ".sql": "SQL query", ".json": "JSON data", ".yaml": "YAML config",
        ".csv": "CSV data", ".md": "Markdown document",
    }.get(ext, "document")

    prompt = (
        f"You are analyzing a {type_hint} file named '{filename}'.\n\n"
        f"Content:\n```\n{truncated}\n```\n\n"
        "Provide a thorough analysis:\n"
        "1. Summarize what this file contains and its purpose.\n"
        "2. Extract all key data, facts, values, and entities.\n"
        "3. Identify any important patterns, issues, or insights.\n"
        "4. For code: explain what it does, highlight key functions/logic.\n"
        "5. Format your response clearly with sections."
    )

    async with httpx.AsyncClient(timeout=120) as client:
        r = await client.post(
            f"{VISION_URL}/api/chat",
            json={
                "model": VISION_MODEL,
                "messages": [{"role": "user", "content": prompt}],
                "stream": False,
            },
        )
        return r.json().get("message", {}).get("content", "").strip()




def _process_pdf(data: bytes) -> tuple[str, int]:
    try:
        import pymupdf4llm
        import fitz  # pymupdf

        doc = fitz.open(stream=data, filetype="pdf")
        pages = doc.page_count

        # Use pymupdf4llm for markdown output (preserves tables, layout)
        md_text = pymupdf4llm.to_markdown(doc)

        if not md_text.strip():
            # Fallback: plain text extraction
            parts = []
            for i, page in enumerate(doc):
                text = page.get_text("text")
                if text.strip():
                    parts.append(f"## Page {i+1}\n\n{text}")
            md_text = "\n\n".join(parts)

        doc.close()
        return md_text, pages
    except Exception as e:
        raise ProcessingError(f"PDF extraction failed: {e}")


# ── DOCX/DOC ──────────────────────────────────────────────────────────────────

def _process_docx(data: bytes, ext: str) -> str:
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            parts = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    style = para.style.name if para.style else ""
                    if "Heading" in style:
                        level = re.search(r"\d+", style)
                        hashes = "#" * (int(level.group()) if level else 1)
                        parts.append(f"{hashes} {text}")
                    else:
                        parts.append(text)
            # Also extract tables
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    parts.append("\n".join(rows))
            return "\n\n".join(parts)
        except Exception as e:
            raise ProcessingError(f"DOCX extraction failed: {e}")
    else:
        # .doc — try extracting as raw text
        try:
            text = _decode_text(data, "file.doc")
            # Clean up binary garbage
            text = re.sub(r"[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f-\x9f]", "", text)
            return text
        except Exception as e:
            raise ProcessingError(f"DOC extraction failed: {e}")


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _process_pptx(data: bytes) -> tuple[str, int]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_parts = [f"## Slide {i+1}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
                # Extract table text
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_parts.append(" | ".join(cells))
            parts.append("\n\n".join(slide_parts))
        return "\n\n---\n\n".join(parts), len(prs.slides)
    except Exception as e:
        raise ProcessingError(f"PPTX extraction failed: {e}")


# ── Excel/ODS ──────────────────────────────────────────────────────────────────

def _process_excel(data: bytes, ext: str) -> tuple[str, int]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                # Skip entirely empty rows
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(rows))
        return "\n\n".join(parts), len(wb.sheetnames)
    except Exception as e:
        raise ProcessingError(f"Excel extraction failed: {e}")


# ── Text / code / structured ──────────────────────────────────────────────────

def _decode_text(data: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()

    # Try UTF-8 first
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        # Use chardet to detect encoding
        detected = chardet.detect(data)
        encoding = detected.get("encoding") or "latin-1"
        try:
            text = data.decode(encoding)
        except Exception:
            text = data.decode("latin-1", errors="replace")

    # For JSON: pretty-print if compact
    if ext in (".json", ".jsonl", ".ndjson"):
        try:
            if ext == ".json":
                obj = json.loads(text)
                text = json.dumps(obj, indent=2, ensure_ascii=False)
            else:
                # JSONL: parse each line
                lines = []
                for line in text.splitlines():
                    line = line.strip()
                    if line:
                        try:
                            lines.append(json.dumps(json.loads(line), indent=2))
                        except Exception:
                            lines.append(line)
                text = "\n\n".join(lines)
        except Exception:
            pass  # Keep as-is

    return text


# ── Helpers ───────────────────────────────────────────────────────────────────

def _is_image(data: bytes) -> bool:
    signatures = [
        b"\xff\xd8\xff",        # JPEG
        b"\x89PNG\r\n\x1a\n",  # PNG
        b"GIF87a", b"GIF89a",   # GIF
        b"RIFF",                # WebP (RIFF....WEBP)
        b"BM",                  # BMP
        b"II*\x00", b"MM\x00*", # TIFF
    ]
    for sig in signatures:
        if data[:len(sig)] == sig:
            return True
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return True
    return False


def _sniff_ext(data: bytes) -> str:
    if _is_image(data):
        if data[:3] == b"\xff\xd8\xff":
            return ".jpg"
        if data[:8] == b"\x89PNG\r\n\x1a\n":
            return ".png"
        return ".jpg"
    if data[:4] == b"%PDF":
        return ".pdf"
    if data[:2] in (b"PK",):
        return ".docx"  # Could be xlsx/pptx too, but default to docx
    return ".txt"


def supported_extensions() -> list[str]:
    return sorted(IMAGE_EXTS | PDF_EXTS | WORD_EXTS | PPT_EXTS | XL_EXTS | TEXT_EXTS)
